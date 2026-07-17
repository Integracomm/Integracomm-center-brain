"""Exportação PPTX do All Hands (pedido Otávio 16/07) — o MESMO conteúdo dos
slides HTML, gerado nativamente com python-pptx: fundo preto #0D0D0D, amarelo
#FFC107, Montserrat (a máquina que abrir precisa ter a fonte; senão o
PowerPoint troca por similar), cards escuros e imagens embutidas.

O layout replica o deck em versão nativa (títulos, cards, barras e tabelas
como formas EDITÁVEIS no PowerPoint) — pequenas diferenças visuais vs o HTML
são esperadas; a versão de referência para projeção continua sendo o
Imprimir/PDF da página gerada.
"""
from __future__ import annotations

import base64
import datetime as dt
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x0D, 0x0D, 0x0D)
CARD = RGBColor(0x1A, 0x1A, 0x1A)
BORDA = RGBColor(0x26, 0x26, 0x26)
AMARELO = RGBColor(0xFF, 0xC1, 0x07)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA = RGBColor(0x9A, 0x9A, 0x9A)
CINZA_ESC = RGBColor(0x77, 0x77, 0x77)
VERMELHO = RGBColor(0xFF, 0x55, 0x55)
PRETO = RGBColor(0x11, 0x11, 0x11)
FONTE = "Montserrat"

_W, _H = Inches(13.333), Inches(7.5)


def _novo_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # em branco
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def _txt(sl, x, y, w, h, linhas, align=PP_ALIGN.LEFT):
    """Caixa de texto; linhas = [(texto, tamanho_pt, cor, negrito, espaco_antes)]."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (texto, tam, cor, negrito, antes) in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if antes:
            p.space_before = Pt(antes)
        r = p.add_run()
        r.text = texto
        f = r.font
        f.name, f.size, f.bold = FONTE, Pt(tam), negrito
        f.color.rgb = cor
    return tb


def _card(sl, x, y, w, h, fill=CARD):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = BORDA
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _cabecalho(sl, kicker, titulo):
    _txt(sl, Inches(0.7), Inches(0.42), Inches(11.9), Inches(1.2), [
        (kicker.upper(), 13, AMARELO, True, 0),
        (titulo.upper(), 27, BRANCO, True, 2),
    ])


def _rodape(sl, rodape):
    _txt(sl, Inches(0.7), Inches(7.06), Inches(8), Inches(0.4),
         [(rodape.upper(), 9, RGBColor(0x55, 0x55, 0x55), False, 0)])


def _num_card(sl, x, y, w, h, valor, rotulo, cor_valor=AMARELO, extra=None):
    _card(sl, x, y, w, h)
    linhas = [(valor, 30, cor_valor, True, 0), (rotulo.upper(), 10, CINZA, False, 4)]
    if extra:
        linhas.append((extra, 8, CINZA_ESC, False, 4))
    _txt(sl, x + Inches(0.15), y + Inches(0.22), w - Inches(0.3), h - Inches(0.4),
         linhas, align=PP_ALIGN.CENTER)


def _tabela(sl, x, y, w, itens, titulo, tot_rotulo="Total"):
    """Card com linhas 'rótulo … n' e total amarelo (substitui a <table> do HTML)."""
    alt = Inches(0.62 + 0.31 * (len(itens) + 1))
    _card(sl, x, y, w, alt)
    linhas = [(titulo.upper(), 10, CINZA, False, 0)]
    for p, n in itens:
        linhas.append((f"{p} — {n}", 12.5, BRANCO, False, 5))
    linhas.append((f"{tot_rotulo} — {sum(n for _, n in itens)}", 13, AMARELO, True, 7))
    _txt(sl, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), alt - Inches(0.3), linhas)
    return y + alt


def _chip_amarelo(sl, x, y, texto, tam=10.5):
    """Chip amarelo compacto (largura pelo texto — como os do deck HTML).
    Retorna o x do próximo chip."""
    w = Inches(0.5 + 0.115 * len(texto))
    c = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.4))
    c.adjustments[0] = 0.22
    c.fill.solid()
    c.fill.fore_color.rgb = AMARELO
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = texto
    r.font.name, r.font.size, r.font.bold = FONTE, Pt(tam), True
    r.font.color.rgb = PRETO
    return x + w + Inches(0.22)


def _foto_b64(data_uri):
    try:
        return io.BytesIO(base64.b64decode(data_uri.split(",", 1)[1]))
    except Exception:  # noqa: BLE001
        return None


def _add_img_ajustada(sl, blob_io, x, y, max_w, max_h):
    """Imagem SEM corte: mantém a proporção dentro da caixa (mesma correção do
    HTML — banner paisagem saía cortado)."""
    from pptx.parts.image import Image as PImage
    blob = blob_io.getvalue()
    try:
        px_w, px_h = PImage.from_blob(blob).size
    except Exception:  # noqa: BLE001
        sl.shapes.add_picture(io.BytesIO(blob), x, y, width=max_w)
        return
    razao = min(max_w / px_w, max_h / px_h)
    w, h = Emu(int(px_w * razao)), Emu(int(px_h * razao))
    sl.shapes.add_picture(io.BytesIO(blob), x + Emu(int((max_w - w) / 2)),
                          y + Emu(int((max_h - h) / 2)), width=w, height=h)


def gerar_pptx(mes: dt.date, d: dict, destaques: list[dict], novos: dict | None,
               orientacoes: list[dict], extras: list[dict]) -> bytes:
    from .allhands import _ASSETS, _MESES, _PLANOS_BARRAS, _fmt_brl0, _plano_barra
    mes_nome = _MESES[mes.month - 1]
    prox = (mes.replace(day=28) + dt.timedelta(days=5)).replace(day=1)
    rodape = f"ALL HANDS · {mes_nome} {mes.year}"
    prs = Presentation()
    prs.slide_width, prs.slide_height = _W, _H

    # 1 — capa (arte oficial de fundo)
    sl = _novo_slide(prs)
    try:
        sl.shapes.add_picture(str(_ASSETS / "allhands_capa.jpg"), 0, 0, width=_W, height=_H)
    except Exception:  # noqa: BLE001 — sem a arte, capa preta
        pass
    _txt(sl, 0, Inches(2.7), _W, Inches(2), [
        ("ALL HANDS", 56, BRANCO, True, 0),
        (f"{mes_nome} / {mes.year}", 18, AMARELO, False, 8),
    ], align=PP_ALIGN.CENTER)

    # 2 — marketing & comercial (funil + meta + vendas por plano em barras)
    sl = _novo_slide(prs)
    _cabecalho(sl, "Marketing & Comercial", "Funil de vendas")
    fun = d["funil"]
    itens_funil = [("👥  TOTAL DE LEADS", f"{fun[0]:,}".replace(",", ".")),
                   ("🎯  OPORTUNIDADES", f"{fun[4]:,}".replace(",", ".")),
                   ("🤝  TOTAL DE VENDAS", f"{d['bookings']:,}".replace(",", "."))]
    y = Inches(1.75)
    for rot, val in itens_funil:
        _card(sl, Inches(0.7), y, Inches(5.2), Inches(0.95))
        _txt(sl, Inches(0.95), y + Inches(0.12), Inches(4.8), Inches(0.75), [
            (rot, 11, BRANCO, True, 0), (val, 20, BRANCO, True, 2)])
        y += Inches(1.13)
    _num_card(sl, Inches(0.7), y + Inches(0.08), Inches(5.2), Inches(1.15),
              _fmt_brl0(d["meta_mes"]), "meta do mês")
    # barras de vendas por plano
    cont = {n: 0 for n in _PLANOS_BARRAS}
    for produto, n, _v in d["vendas_plano"]:
        cont[_plano_barra(produto)] += n
    vmax = max(cont.values()) or 1
    _txt(sl, Inches(6.6), Inches(1.6), Inches(6), Inches(0.4),
         [("VENDAS POR PLANO", 15, AMARELO, True, 0)])
    _card(sl, Inches(6.6), Inches(2.05), Inches(6.05), Inches(4.6))
    yb = Inches(2.35)
    for nome in _PLANOS_BARRAS:
        v = cont[nome]
        _txt(sl, Inches(6.8), yb, Inches(2.15), Inches(0.35), [(nome, 10.5, BRANCO, True, 0)])
        larg = Inches(0.05 + 3.1 * v / vmax)
        if v:
            bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), yb + Inches(0.02),
                                      larg, Inches(0.28))
            bar.adjustments[0] = 0.3
            bar.fill.solid()
            bar.fill.fore_color.rgb = AMARELO
            bar.line.fill.background()
            bar.shadow.inherit = False
        _txt(sl, Inches(9.05) + larg, yb - Inches(0.02), Inches(0.8), Inches(0.35),
             [(str(v), 12, BRANCO, True, 0)])
        yb += Inches(0.72)
    _rodape(sl, rodape)

    # 3 — clientes por plano (antigos × novos)
    from .allhands import _PLANOS_NOVOS
    sl = _novo_slide(prs)
    _cabecalho(sl, "Operação", "Clientes por plano")
    _txt(sl, Inches(0.7), Inches(1.55), Inches(11), Inches(0.7), [
        (f"{d['base_ativa']}", 38, AMARELO, True, 0)])
    _txt(sl, Inches(2.0), Inches(1.78), Inches(8), Inches(0.4), [
        ("TOTAL GERAL DE CLIENTES ATIVOS", 11, CINZA, False, 0)])
    col_novos = [(p, n) for p, n in d["clientes_plano"] if any(k in p.lower() for k in _PLANOS_NOVOS)]
    col_ant = [(p, n) for p, n in d["clientes_plano"] if (p, n) not in col_novos]
    _tabela(sl, Inches(0.7), Inches(2.5), Inches(5.9), col_ant[:9], "Planos antigos")
    _tabela(sl, Inches(6.85), Inches(2.5), Inches(5.9), col_novos[:9], "Planos novos")
    _rodape(sl, rodape)

    # 4 — saídas (churn) — régua recorrente (sem B1) + MRR perdido
    sl = _novo_slide(prs)
    _cabecalho(sl, "Retenção", "Saídas de clientes (churn)")
    sd_novos = [(p, n) for p, n in d["saidas_plano"] if any(k in p.lower() for k in _PLANOS_NOVOS)]
    sd_ant = [(p, n) for p, n in d["saidas_plano"] if (p, n) not in sd_novos]
    _tabela(sl, Inches(0.7), Inches(1.8), Inches(3.9), sd_ant[:7], "Planos antigos")
    _tabela(sl, Inches(4.8), Inches(1.8), Inches(3.9), sd_novos[:7], "Planos novos")
    taxa_txt = (f"{d['taxa_mes'] * 100:.1f}%".replace(".", ",")
                if d.get("taxa_mes") is not None else "—")
    _num_card(sl, Inches(8.9), Inches(1.8), Inches(3.75), Inches(2.0), taxa_txt,
              "taxa de cancelamento do mês", cor_valor=VERMELHO,
              extra=(f"{_fmt_brl0(d.get('mrr_rec_perdido'))} de MRR perdido · cancelados recorrentes ÷ "
                     f"clientes recorrentes ativos ({d.get('saidas_rec', 0)} de {d.get('base_rec', 0)}; "
                     f"cada cliente conta 1×) — B1/Start fora"))
    meta_ch = f"{(d.get('churn_meta') or 0.05) * 100:.0f}%"
    _num_card(sl, Inches(8.9), Inches(4.0), Inches(3.75), Inches(1.5), meta_ch,
              "meta taxa de cancelamento")
    _rodape(sl, rodape)

    # 5 — evolução da taxa de cancelamento (barras)
    sl = _novo_slide(prs)
    _cabecalho(sl, "Retenção", "Evolução da taxa de cancelamento")
    serie = d.get("churn_serie") or []
    if serie:
        vmax = max(v for _, v in serie) or 1
        n = len(serie)
        larg = Inches(0.9)
        gap = (Inches(11.9) - larg * n) / max(1, n - 1) if n > 1 else Inches(0)
        x = Inches(0.7)
        for nome, v in serie:
            h_bar = Emu(int(Inches(3.4) * (v / vmax))) if vmax else Inches(0.1)
            y_bar = Inches(5.9) - h_bar
            bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_bar, larg, h_bar)
            bar.adjustments[0] = 0.12
            bar.fill.solid()
            bar.fill.fore_color.rgb = AMARELO
            bar.line.fill.background()
            bar.shadow.inherit = False
            _txt(sl, x - Inches(0.25), y_bar - Inches(0.42), larg + Inches(0.5), Inches(0.4),
                 [(f"{v * 100:.1f}%".replace(".", ","), 13, BRANCO, True, 0)], align=PP_ALIGN.CENTER)
            _txt(sl, x - Inches(0.25), Inches(6.0), larg + Inches(0.5), Inches(0.4),
                 [(nome, 11, CINZA, False, 0)], align=PP_ALIGN.CENTER)
            x += larg + gap
        _txt(sl, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4), [
            ("série = realizado lançado na planilha de planejamento (régua da gestão)",
             9, CINZA_ESC, False, 0)], align=PP_ALIGN.CENTER)
    _rodape(sl, rodape)

    # 6 — metas do próximo mês
    sl = _novo_slide(prs)
    _cabecalho(sl, f"Metas {_MESES[prox.month - 1].title()}", "Fechamentos & Retenção")
    churn_meta_txt = (f"{d['churn_meta'] * 100:.0f}%" if d.get("churn_meta") is not None else "5%")
    _num_card(sl, Inches(2.2), Inches(2.6), Inches(4.2), Inches(2.1),
              _fmt_brl0(d.get("meta_prox")), "meta do mês — comercial")
    _num_card(sl, Inches(6.9), Inches(2.6), Inches(4.2), Inches(2.1),
              churn_meta_txt, "meta taxa de cancelamento")
    _rodape(sl, rodape)

    # 7/8 — destaques & promoções (manuais) — intro no layout do HTML (16/07):
    # estrela + kicker + título + texto em bloco centralizado na vertical,
    # chips compactos logo abaixo (não os blocões soltos da 1ª versão)
    if destaques:
        sl = _novo_slide(prs)
        try:
            sl.shapes.add_picture(str(_ASSETS / "allhands_estrela.png"),
                                  Inches(0.7), Inches(1.5), width=Inches(0.72))
        except Exception:  # noqa: BLE001 — sem a arte, segue sem estrela
            pass
        _txt(sl, Inches(0.7), Inches(2.45), Inches(9.6), Inches(2.4), [
            ("PROMOVIDOS & DESTAQUES DO MÊS", 13, AMARELO, True, 0),
            ("CELEBRANDO CONQUISTAS", 34, BRANCO, True, 4),
            ("Um momento para parabenizar aqueles que demonstraram evolução e elevaram "
             "o padrão de excelência da nossa equipe.", 14, CINZA, False, 10)])
        xc = Inches(0.7)
        for chip in ("DEDICAÇÃO", "RESULTADOS", "CULTURA"):
            xc = _chip_amarelo(sl, xc, Inches(5.1), chip)
        _rodape(sl, rodape)
        for dq in destaques:
            sl = _novo_slide(prs)
            foto = _foto_b64(dq["foto"]) if dq.get("foto") else None
            if foto:
                _add_img_ajustada(sl, foto, Inches(0.9), Inches(1.4), Inches(3.4), Inches(4.6))
            else:
                _card(sl, Inches(0.9), Inches(1.4), Inches(3.4), Inches(4.6))
            _txt(sl, Inches(5.0), Inches(1.75), Inches(7.4), Inches(0.4),
                 [("DESTAQUES & PROMOÇÕES", 12, AMARELO, True, 0)])
            _chip_amarelo(sl, Inches(5.0), Inches(2.3), dq["tipo"].upper(), tam=11)
            _txt(sl, Inches(5.0), Inches(2.95), Inches(7.4), Inches(2.6), [
                (dq["nome"], 34, BRANCO, True, 0),
                (dq["cargo"].upper(), 14, AMARELO, False, 8),
            ] + ([(dq["nivel"], 12, CINZA, False, 6)] if dq.get("nivel") else []))
            _rodape(sl, rodape)

    # 9 — novos colaboradores
    if novos:
        sl = _novo_slide(prs)
        _cabecalho(sl, "Apresentação de novos colaboradores", "Sejam bem-vindos ao time!")
        if novos.get("nomes"):
            _txt(sl, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
                 [(novos["nomes"], 14, AMARELO, False, 0)])
        x = Inches(0.7)
        for ic, t, q in (("🪪", "NOME + EQUIPE", "Qual seu nome e equipe que está trabalhando?"),
                         ("💬", "PRIMEIROS DIAS", "Como está sendo seus primeiros dias aqui conosco?"),
                         ("🚀", "HOBBIES", "O que gosta de fazer nas horas vagas?")):
            _card(sl, x, Inches(2.6), Inches(3.9), Inches(2.6))
            _txt(sl, x + Inches(0.3), Inches(2.9), Inches(3.3), Inches(2.1), [
                (f"{ic}  {t}", 14, BRANCO, True, 0), (q, 12, CINZA, False, 10)])
            x += Inches(4.1)
        _rodape(sl, rodape)

    # 10 — orientações & novidades — layout do HTML (16/07): cards CENTRADOS
    # na página (não encostados no canto), ícone em círculo escuro no topo
    if orientacoes:
        sl = _novo_slide(prs)
        _cabecalho(sl, "Convivência e organização do espaço", "Orientações & Novidades")
        os_ = orientacoes[:6]
        por_linha = min(3, len(os_))
        n_linhas = (len(os_) + por_linha - 1) // por_linha
        card_w = Inches(4.5 if len(os_) <= 2 else 3.9)
        card_h = Inches(3.2 if n_linhas == 1 else 2.55)
        gap = Inches(0.3)
        area_top, area_bot = Inches(1.65), Inches(6.95)
        y0 = area_top + Emu(int((area_bot - area_top - (n_linhas * card_h + (n_linhas - 1) * gap)) / 2))
        for i, o in enumerate(os_):
            lin, col = divmod(i, por_linha)
            n_nesta = min(por_linha, len(os_) - lin * por_linha)
            total_w = n_nesta * card_w + (n_nesta - 1) * gap
            x = Emu(int((_W - total_w) / 2)) + col * (card_w + gap)
            y = y0 + lin * (card_h + gap)
            _card(sl, x, y, card_w, card_h)
            dia = Inches(1.05)
            ov = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Emu(int((card_w - dia) / 2)),
                                     y + Inches(0.35), dia, dia)
            ov.fill.solid()
            ov.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
            ov.line.fill.background()
            ov.shadow.inherit = False
            tf = ov.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = o["icone"]
            r.font.size = Pt(30)
            _txt(sl, x + Inches(0.25), y + Inches(1.6), card_w - Inches(0.5), card_h - Inches(1.7), [
                (o["titulo"].upper(), 15, BRANCO, True, 0),
            ] + ([(o["texto"], 11.5, CINZA, False, 7)] if o.get("texto") else []),
                 align=PP_ALIGN.CENTER)
        _rodape(sl, rodape)

    # extras — slides livres (cada linha = marcador; imagem sem corte)
    for ex in (extras or []):
        sl = _novo_slide(prs)
        _cabecalho(sl, ex.get("kicker") or "Novidades", ex["titulo"])
        linhas_ex = [ln.strip() for ln in (ex.get("texto") or "").splitlines() if ln.strip()]
        tem_img = bool(ex.get("foto"))
        if linhas_ex:
            larg = Inches(6.4) if tem_img else Inches(11.9)
            _card(sl, Inches(0.7), Inches(2.1), larg, Inches(0.6 + 0.52 * len(linhas_ex)))
            _txt(sl, Inches(1.0), Inches(2.35), larg - Inches(0.6), Inches(0.4 + 0.52 * len(linhas_ex)),
                 [(f"●  {ln}", 14, BRANCO, False, 10 if i else 0) for i, ln in enumerate(linhas_ex)])
        if tem_img:
            foto = _foto_b64(ex["foto"])
            if foto:
                _add_img_ajustada(sl, foto, Inches(7.5), Inches(1.7), Inches(5.1), Inches(4.9))
        _rodape(sl, rodape)

    # fechamento
    sl = _novo_slide(prs)
    _txt(sl, 0, Inches(2.8), _W, Inches(2), [
        ("ÓTIMO MÊS", 42, BRANCO, True, 0),
        ("PARA TODOS NÓS", 42, AMARELO, True, 2),
        (rodape.upper(), 11, CINZA, False, 14),
    ], align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
